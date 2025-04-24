from flask import Flask, request, send_file, jsonify
from flask_cors import CORS
import xlwings as xw
from pptx import Presentation
from io import BytesIO
import re

app = Flask(__name__)
CORS(app)

@app.route('/')
def home():
    return '''
    <h1>👋 Welcome to Project Senior</h1>
    <p>The backend is running successfully.</p>
    '''


def normalize(text):
    return re.sub(r"\s+", " ", text.strip().lower())

def is_exec_summary_table(table):
    headers = [normalize(cell.text) for cell in table.rows[0].cells]
    expected = ["", "impressions", "unique impressions", "spend", "attributed sku revenue", "sku roas"]
    return headers[:6] == expected

def extract_exec_summary_data(prs):
    for slide in prs.slides:
        if "offsite campaign executive summary" in normalize(" ".join([s.text for s in slide.shapes if hasattr(s, 'text')])):
            for shape in slide.shapes:
                if shape.has_table and is_exec_summary_table(shape.table):
                    table = shape.table
                    data_rows = []
                    for i in range(1, len(table.rows) - 1):
                        data_rows.append([cell.text.strip() for cell in table.rows[i].cells])
                    return data_rows
    return []

def apply_campaign_filters(ws, campaign_ids, group_type):
    row_map = [6, 40, 74, 111, 148]
    for idx, campaign in enumerate(campaign_ids[:5]):
        if campaign:
            cell = ws.range(f"B{row_map[idx]}")
            cell.value = campaign
    if group_type == 'Pacing':
        expected = "All Line Items For Campaign"
        for r in row_map:
            if ws.range(f"B{r}").value.strip() != expected:
                raise ValueError(f"Expected '{expected}' in B{r}, found '{ws.range(f'B{r}').value}'")
    elif group_type == 'IO':
        if campaign_ids[4] or campaign_ids[5]:
            combo = list(filter(None, [campaign_ids[4], campaign_ids[5]]))
            ws.range("B148").value = ", ".join(combo)

def update_exec_summary(ws, data_rows, start_row):
    for i, row in enumerate(data_rows):
        for j, val in enumerate(row):
            ws.range((start_row + i, j + 1)).value = val

def is_green(rgb):
    return rgb[1] > 180 and rgb[0] < 150 and rgb[2] < 150

def is_red(rgb):
    return rgb[0] > 180 and rgb[1] < 150 and rgb[2] < 150

def check_color_status(ws):
    row_blocks = [(10, 34), (44, 68), (78, 102), (115, 139), (152, 176)]
    for start, end in row_blocks:
        for row in range(start, end + 1):
            e_color = ws.range(f"E{row}").color
            f_color = ws.range(f"F{row}").color
            if e_color and f_color:
                if is_red(e_color) and is_red(f_color):
                    return "Complete Failure"
                if (is_green(e_color) and is_green(f_color)) or (is_green(e_color) and is_red(f_color)) or (is_red(e_color) and is_green(f_color)):
                    return "Complete Success"
    return "No matching condition"

@app.route('/api/process', methods=['POST'])
def process():
    try:
        ppt = request.files.get('ppt')
        excel = request.files.get('excel')
        vba = request.files.get('vba')

        row_data = {f"row{i}": request.form.get(f"row{i}") for i in range(9, 13)}
        campaign_type = request.form.get("campaignType")
        group_type = request.form.get("groupType")
        campaign_ids = [request.form.get(f"ioInput{i+1}") for i in range(6)]

        temp_ppt_path = "/tmp/temp.pptx"
        temp_excel_path = "/tmp/temp.xlsm"
        ppt.save(temp_ppt_path)
        excel.save(temp_excel_path)

        app_xl = xw.App(visible=False)
        wb = app_xl.books.open(temp_excel_path)
        ws1 = wb.sheets["Sheet1"]
        ws2 = wb.sheets["COPY PASTE"]

        for k, v in row_data.items():
            if v:
                row_num = int(k.replace("row", ""))
                ws2.range(f"B{row_num}").value = v

        prs = Presentation(temp_ppt_path)
        exec_data = extract_exec_summary_data(prs)
        if campaign_type == 'SB':
            update_exec_summary(ws2, exec_data, 32)
        elif campaign_type == 'MBC':
            update_exec_summary(ws2, exec_data, 42)

        apply_campaign_filters(ws1, campaign_ids, group_type)
        color_status = check_color_status(ws1)

        output = BytesIO()
        wb.save("/tmp/final.xlsm")
        wb.close()
        app_xl.quit()

        with open("/tmp/final.xlsm", "rb") as f:
            output.write(f.read())
        output.seek(0)
        response = send_file(output, download_name="Updated_Report.xlsm", as_attachment=True)
        response.headers["X-Color-Status"] = color_status
        return response

    except Exception as e:
        return jsonify({"error": str(e)}), 500

if __name__ == "__main__":
    app.run(host="0.0.0.0", port=10000)